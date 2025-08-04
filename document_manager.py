import os
import chromadb
from sentence_transformers import SentenceTransformer
import hashlib
import warnings
from config import load_settings, get_docs_path

# 경고 메시지 숨기기
warnings.filterwarnings("ignore", category=FutureWarning)
warnings.filterwarnings("ignore", message=".*encoder_attention_mask.*")
warnings.filterwarnings("ignore", message=".*torch.*")

# 다양한 파일 형식 처리를 위한 라이브러리들
try:
    from docx import Document  # Word 파일
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2  # PDF 파일
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import openpyxl  # Excel 파일
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation  # PowerPoint 파일
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import fitz  # PyMuPDF (더 나은 PDF 처리)
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from markitdown import MarkItDown  # Microsoft의 통합 문서 변환 라이브러리
    MARKITDOWN_AVAILABLE = True
except ImportError:
    MARKITDOWN_AVAILABLE = False

print("=== TinyRAG - 가벼운 문서 검색 시스템 ===")

class DocumentProcessor:
    def __init__(self, settings_file="settings.json"):
        # 설정 로드
        self.settings = load_settings(settings_file)
        
        self.db_path = self.settings["paths"]["chroma_db"]
        self.model_path = self.settings["embedding_model"]["local_path"]
        
        # ChromaDB 클라이언트 생성
        print("ChromaDB 클라이언트 초기화 중...")
        self.client = chromadb.PersistentClient(path=self.db_path)
        
        # SentenceTransformer 모델 로딩
        print("임베딩 모델 로딩 중...")
        self.model = SentenceTransformer(self.model_path)
        
        # MarkItDown 초기화 (사용 가능한 경우)
        if MARKITDOWN_AVAILABLE:
            print("MarkItDown 문서 변환기 초기화 중...")
            self.markitdown = MarkItDown()
        else:
            self.markitdown = None
            
        print("✅ 시스템 초기화 완료")
    
    def split_text(self, text, chunk_size=None, overlap=None):
        """텍스트를 자연스러운 구분점에서 청크로 분할"""
        if chunk_size is None:
            chunk_size = self.settings["search"]["chunk_size"]
        if overlap is None:
            overlap = self.settings["search"]["overlap"]
            
        chunks = []
        metadata = []
        start = 0
        chunk_id = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # 텍스트 끝에 도달한 경우
            if end >= len(text):
                chunk_text = text[start:].strip()
            else:
                # 자연스러운 구분점에서 자르기
                potential_chunk = text[start:end]
                
                # 1. 문장의 끝(.)에서 자르기 시도
                last_period = potential_chunk.rfind('.')
                if last_period > chunk_size * 0.7:  # 너무 짧지 않다면
                    actual_end = start + last_period + 1
                    chunk_text = text[start:actual_end].strip()
                else:
                    # 2. 줄바꿈에서 자르기 시도
                    last_newline = potential_chunk.rfind('\n')
                    if last_newline > chunk_size * 0.7:
                        actual_end = start + last_newline
                        chunk_text = text[start:actual_end].strip()
                    else:
                        # 3. 공백에서 자르기 시도
                        last_space = potential_chunk.rfind(' ')
                        if last_space > chunk_size * 0.7:
                            actual_end = start + last_space
                            chunk_text = text[start:actual_end].strip()
                        else:
                            # 4. 마지막 수단으로 원래 크기에서 자르기
                            chunk_text = potential_chunk.strip()
                            actual_end = end
            
            if chunk_text:  # 빈 청크 제외
                chunks.append(chunk_text)
                metadata.append({
                    "chunk_id": chunk_id,
                    "start_pos": start,
                    "end_pos": min(actual_end if 'actual_end' in locals() else end, len(text)),
                    "length": len(chunk_text)
                })
                chunk_id += 1
            
            # 다음 청크 시작점 계산 (겹침 고려)
            if 'actual_end' in locals():
                start = max(actual_end - overlap, start + 1)
                del actual_end  # 변수 초기화
            else:
                start += chunk_size - overlap
        
        return chunks, metadata
    
    def get_supported_extensions(self):
        """지원되는 파일 확장자 목록 반환"""
        extensions = ['.txt', '.md']  # 기본 텍스트 파일
        
        if DOCX_AVAILABLE:
            extensions.extend(['.docx', '.doc'])
        if PDF_AVAILABLE or PYMUPDF_AVAILABLE:
            extensions.append('.pdf')
        if EXCEL_AVAILABLE:
            extensions.extend(['.xlsx', '.xls'])
        if PPTX_AVAILABLE:
            extensions.extend(['.pptx', '.ppt'])
            
        return extensions
    
    def read_word_file(self, file_path):
        """Word 파일을 Markdown 형식으로 읽기 - MarkItDown 우선 사용"""
        # MarkItDown 우선 시도
        if MARKITDOWN_AVAILABLE:
            content = self.read_office_file_with_markitdown(file_path)
            if content is not None:
                return content
        
        # MarkItDown 실패시 기존 방법 사용
        if not DOCX_AVAILABLE:
            raise ValueError("python-docx 라이브러리가 설치되지 않았습니다.")
        
        try:
            doc = Document(file_path)
            markdown_content = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                    
                # 제목 스타일 처리
                if paragraph.style.name.startswith('Heading'):
                    level = paragraph.style.name.replace('Heading ', '')
                    if level.isdigit():
                        markdown_content.append(f"{'#' * int(level)} {text}")
                    else:
                        markdown_content.append(f"## {text}")
                else:
                    markdown_content.append(text)
            
            # 표 처리
            for table in doc.tables:
                markdown_content.append("\n| " + " | ".join([cell.text for cell in table.rows[0].cells]) + " |")
                markdown_content.append("| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |")
                for row in table.rows[1:]:
                    markdown_content.append("| " + " | ".join([cell.text for cell in row.cells]) + " |")
                markdown_content.append("")
            
            return "\n\n".join(markdown_content)
        except Exception as e:
            raise ValueError(f"Word 파일 읽기 오류: {e}")
    
    def read_pdf_file(self, file_path):
        """PDF 파일을 텍스트로 읽기 - MarkItDown 우선 사용"""
        if MARKITDOWN_AVAILABLE:
            return self._read_pdf_markitdown(file_path)
        elif PYMUPDF_AVAILABLE:
            return self._read_pdf_pymupdf(file_path)
        elif PDF_AVAILABLE:
            return self._read_pdf_pypdf2(file_path)
        else:
            raise ValueError("PDF 처리 라이브러리가 설치되지 않았습니다.")
    
    def _read_pdf_markitdown(self, file_path):
        """MarkItDown을 사용한 PDF 읽기 (최고 품질)"""
        try:
            print("MarkItDown을 사용하여 PDF 변환 중...")
            result = self.markitdown.convert(file_path)
            return result.text_content
        except Exception as e:
            print(f"⚠️ MarkItDown 변환 실패, PyMuPDF로 대체 시도: {e}")
            # MarkItDown 실패시 다른 방법으로 대체
            if PYMUPDF_AVAILABLE:
                return self._read_pdf_pymupdf(file_path)
            elif PDF_AVAILABLE:
                return self._read_pdf_pypdf2(file_path)
            else:
                raise ValueError(f"PDF 파일 읽기 오류 (MarkItDown): {e}")
    
    def _read_pdf_pymupdf(self, file_path):
        """PyMuPDF를 사용한 PDF 읽기 (더 정확함)"""
        try:
            doc = fitz.open(file_path)
            text_content = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    text_content.append(f"## 페이지 {page_num + 1}\n\n{text}")
            
            doc.close()
            return "\n\n".join(text_content)
        except Exception as e:
            raise ValueError(f"PDF 파일 읽기 오류 (PyMuPDF): {e}")
    
    def _read_pdf_pypdf2(self, file_path):
        """PyPDF2를 사용한 PDF 읽기"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_content = []
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(f"## 페이지 {page_num + 1}\n\n{text}")
                
                return "\n\n".join(text_content)
        except Exception as e:
            raise ValueError(f"PDF 파일 읽기 오류 (PyPDF2): {e}")
    
    def read_excel_file(self, file_path):
        """Excel 파일을 Markdown 표 형식으로 읽기 - MarkItDown 우선 사용"""
        # MarkItDown 우선 시도
        if MARKITDOWN_AVAILABLE:
            content = self.read_office_file_with_markitdown(file_path)
            if content is not None:
                return content
        
        # MarkItDown 실패시 기존 방법 사용
        if not EXCEL_AVAILABLE:
            raise ValueError("openpyxl 라이브러리가 설치되지 않았습니다.")
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            markdown_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                markdown_content.append(f"# 시트: {sheet_name}\n")
                
                # 데이터가 있는 영역 찾기
                max_row = sheet.max_row
                max_col = sheet.max_column
                
                if max_row > 0 and max_col > 0:
                    # 헤더 행
                    header_row = []
                    for col in range(1, max_col + 1):
                        cell_value = sheet.cell(row=1, column=col).value
                        header_row.append(str(cell_value) if cell_value is not None else "")
                    
                    markdown_content.append("| " + " | ".join(header_row) + " |")
                    markdown_content.append("| " + " | ".join(["---"] * len(header_row)) + " |")
                    
                    # 데이터 행들
                    for row in range(2, min(max_row + 1, 101)):  # 최대 100행까지만
                        data_row = []
                        for col in range(1, max_col + 1):
                            cell_value = sheet.cell(row=row, column=col).value
                            data_row.append(str(cell_value) if cell_value is not None else "")
                        markdown_content.append("| " + " | ".join(data_row) + " |")
                
                markdown_content.append("")
            
            return "\n\n".join(markdown_content)
        except Exception as e:
            raise ValueError(f"Excel 파일 읽기 오류: {e}")
    
    def read_office_file_with_markitdown(self, file_path):
        """MarkItDown을 사용하여 Office 문서 읽기 (Word, Excel, PowerPoint)"""
        if not MARKITDOWN_AVAILABLE:
            return None
            
        try:
            print("MarkItDown을 사용하여 Office 문서 변환 중...")
            result = self.markitdown.convert(file_path)
            return result.text_content
        except Exception as e:
            print(f"⚠️ MarkItDown Office 문서 변환 실패: {e}")
            return None
    
    def read_powerpoint_file(self, file_path):
        """PowerPoint 파일을 Markdown 형식으로 읽기 - MarkItDown 우선 사용"""
        # MarkItDown 우선 시도
        if MARKITDOWN_AVAILABLE:
            content = self.read_office_file_with_markitdown(file_path)
            if content is not None:
                return content
        
        # MarkItDown 실패시 기존 방법 사용
        if not PPTX_AVAILABLE:
            raise ValueError("python-pptx 라이브러리가 설치되지 않았습니다.")
        
        try:
            prs = Presentation(file_path)
            markdown_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_content.append(f"# 슬라이드 {slide_num}\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # 제목인 경우 (일반적으로 첫 번째 텍스트 박스)
                        if shape == slide.shapes[0]:
                            markdown_content.append(f"## {text}\n")
                        else:
                            markdown_content.append(f"{text}\n")
                
                markdown_content.append("")
            
            return "\n\n".join(markdown_content)
        except Exception as e:
            raise ValueError(f"PowerPoint 파일 읽기 오류: {e}")
    
    def detect_file_type_and_read(self, file_path):
        """파일 확장자를 감지하고 적절한 방법으로 읽기"""
        _, ext = os.path.splitext(file_path.lower())
        
        print(f"파일 형식 감지: {ext}")
        
        if ext in ['.txt', '.md']:
            return self.load_document_text(file_path)
        elif ext in ['.docx', '.doc']:
            content = self.read_word_file(file_path)
            return content, 'utf-8'
        elif ext == '.pdf':
            content = self.read_pdf_file(file_path)
            return content, 'utf-8'
        elif ext in ['.xlsx', '.xls']:
            content = self.read_excel_file(file_path)
            return content, 'utf-8'
        elif ext in ['.pptx', '.ppt']:
            content = self.read_powerpoint_file(file_path)
            return content, 'utf-8'
        else:
            # 알 수 없는 확장자는 텍스트로 시도
            print(f"⚠️ 알 수 없는 파일 형식입니다. 텍스트 파일로 시도합니다: {ext}")
            return self.load_document_text(file_path)
    
    def load_document_text(self, file_path, encoding_list=['euc-kr', 'utf-8', 'cp949']):
        """텍스트 문서 로드 (여러 인코딩 시도)"""
        print(f"텍스트 문서 로딩 중: {file_path}")
        
        for encoding in encoding_list:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                print(f"✅ {encoding} 인코딩으로 문서 로드 성공")
                return text, encoding
            except UnicodeDecodeError:
                continue
        
        raise ValueError(f"지원하는 인코딩으로 파일을 읽을 수 없습니다: {encoding_list}")
    
    def load_document(self, file_path):
        """통합 문서 로더 - 파일 형식에 따라 자동으로 적절한 방법 선택"""
        return self.detect_file_type_and_read(file_path)
        return self.detect_file_type_and_read(file_path)
    
    def create_collection(self, collection_name, document_path):
        """문서로부터 ChromaDB 컬렉션 생성"""
        print(f"\n=== '{collection_name}' 컬렉션 생성 ===")
        
        # 컬렉션이 이미 존재하는지 확인
        existing_collections = [col.name for col in self.client.list_collections()]
        if collection_name in existing_collections:
            print(f"⚠️  '{collection_name}' 컬렉션이 이미 존재합니다.")
            response = input("기존 컬렉션을 덮어쓰시겠습니까? (y/N): ").strip().lower()
            if response not in ['y', 'yes']:
                print("작업을 취소했습니다.")
                return None
            else:
                self.client.delete_collection(name=collection_name)
                print("기존 컬렉션을 삭제했습니다.")
        
        # 새 컬렉션 생성
        collection = self.client.create_collection(name=collection_name)
        
        # 문서 로드
        text, encoding = self.load_document(document_path)
        print(f"문서 크기: {len(text)} 문자")
        
        # 텍스트 분할
        # 큰 문서의 경우 청크 크기 자동 조정
        large_doc_threshold = self.settings["search"]["large_doc_threshold"]
        if len(text) > large_doc_threshold:  # 설정값 이상인 경우
            chunk_size = self.settings["search"]["large_doc_chunk_size"]
            overlap = self.settings["search"]["large_doc_overlap"]
            print(f"📄 큰 문서 감지 - 청크 크기를 {chunk_size}자로 조정합니다.")
        else:
            chunk_size = self.settings["search"]["chunk_size"]
            overlap = self.settings["search"]["overlap"]
            
        chunks, metadata = self.split_text(text, chunk_size=chunk_size, overlap=overlap)
        print(f"텍스트를 {len(chunks)}개 청크로 분할했습니다.")
        
        # ChromaDB에 저장하면서 개별적으로 임베딩 생성
        print("청크별 임베딩 생성 및 저장 중...")
        
        # 문서 ID 생성 (파일명 기반)
        doc_name = os.path.basename(document_path)
        doc_id_base = doc_name.replace('.', '_')
        
        # 청크를 개별적으로 처리하여 ChromaDB에 저장
        batch_size = self.settings["embedding_model"].get("batch_size", 1000)  # 설정에서 가져오거나 기본값 사용
        total_processed = 0
        
        print(f"배치 크기: {batch_size}개씩 처리합니다.")
        
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i:i+batch_size]
            batch_metadata = metadata[i:i+batch_size]
            current_batch_size = len(batch_chunks)
            
            print(f"   배치 {i//batch_size + 1}/{(len(chunks) + batch_size - 1)//batch_size} 처리 중... ({current_batch_size}개 청크)")
            
            try:
                # 배치별 임베딩 생성
                batch_embeddings = self.model.encode(batch_chunks, show_progress_bar=False)
                
                # 배치별 ID 및 메타데이터 생성
                batch_ids = [f"{doc_id_base}_chunk_{i+j}" for j in range(current_batch_size)]
                
                # 메타데이터에 문서 정보 추가
                enhanced_batch_metadata = []
                for j, meta in enumerate(batch_metadata):
                    meta.update({
                        "document_name": doc_name,
                        "document_path": document_path,
                        "encoding": encoding,
                        "chunk_text_preview": batch_chunks[j][:100] + "..." if len(batch_chunks[j]) > 100 else batch_chunks[j]
                    })
                    enhanced_batch_metadata.append(meta)
                
                # ChromaDB에 배치 추가
                collection.add(
                    embeddings=[emb.tolist() if hasattr(emb, 'tolist') else emb for emb in batch_embeddings],
                    documents=batch_chunks,
                    metadatas=enhanced_batch_metadata,
                    ids=batch_ids
                )
                
                total_processed += current_batch_size
                print(f"     ✅ {current_batch_size}개 청크 저장 완료 (총 {total_processed}/{len(chunks)})")
                
            except Exception as e:
                print(f"     ❌ 배치 {i//batch_size + 1} 처리 중 오류: {e}")
                
                # 개별 청크로 다시 시도
                print(f"     🔄 개별 청크 처리로 재시도...")
                for j, chunk in enumerate(batch_chunks):
                    try:
                        # 개별 임베딩 생성
                        embedding = self.model.encode([chunk], show_progress_bar=False)[0]
                        
                        # 개별 ID 및 메타데이터
                        chunk_id = f"{doc_id_base}_chunk_{i+j}"
                        chunk_meta = batch_metadata[j].copy()
                        chunk_meta.update({
                            "document_name": doc_name,
                            "document_path": document_path,
                            "encoding": encoding,
                            "chunk_text_preview": chunk[:100] + "..." if len(chunk) > 100 else chunk
                        })
                        
                        # ChromaDB에 개별 추가
                        collection.add(
                            embeddings=[embedding.tolist() if hasattr(embedding, 'tolist') else embedding],
                            documents=[chunk],
                            metadatas=[chunk_meta],
                            ids=[chunk_id]
                        )
                        
                        total_processed += 1
                        
                    except Exception as inner_e:
                        print(f"       ⚠️ 청크 {i+j+1} 스킵 (오류: {inner_e})")
                        continue
        
        print(f"✅ 총 {total_processed}개 청크 처리 완료")
        
        # ChromaDB에 저장 완료 메시지 (기존 저장 로직 제거)
        
        print(f"✅ '{collection_name}' 컬렉션 생성 완료!")
        print(f"   - 문서: {doc_name}")
        print(f"   - 처리된 청크 수: {total_processed}")
        print(f"   - 인코딩: {encoding}")
        
        return collection
    
    def list_collections(self):
        """저장된 컬렉션 목록 출력"""
        collections = self.client.list_collections()
        print(f"\n=== 저장된 컬렉션 목록 ===")
        if not collections:
            print("저장된 컬렉션이 없습니다.")
            return collections
            
        for i, collection in enumerate(collections):
            print(f"{i+1}. {collection.name}")
            # 컬렉션 정보 가져오기
            count = collection.count()
            print(f"   - 문서 청크 수: {count}")
            
            # 컬렉션의 첫 번째 문서 정보 가져오기
            try:
                sample = collection.peek(limit=1)
                if sample['metadatas'] and len(sample['metadatas']) > 0:
                    metadata = sample['metadatas'][0]
                    doc_name = metadata.get('document_name', 'Unknown')
                    doc_path = metadata.get('document_path', 'Unknown')
                    encoding = metadata.get('encoding', 'Unknown')
                    print(f"   - 문서명: {doc_name}")
                    print(f"   - 파일경로: {doc_path}")
                    print(f"   - 인코딩: {encoding}")
            except:
                print("   - 메타데이터 정보를 가져올 수 없습니다.")
            print()
        
        return collections
    
    def generate_collection_name(self, document_path):
        """문서 경로를 기반으로 고유한 컬렉션명 생성"""
        doc_name = os.path.basename(document_path)
        name_without_ext = os.path.splitext(doc_name)[0]
        # 특수문자를 언더스코어로 변경
        safe_name = "".join(c if c.isalnum() or c in '-_' else '_' for c in name_without_ext)
        return safe_name
    
    def add_document(self, document_path, collection_name=None):
        """문서를 고유한 컬렉션으로 추가"""
        # docs 폴더에서 파일 찾기 시도
        full_path = get_docs_path(document_path)
        
        if not os.path.exists(full_path):
            print(f"❌ 파일을 찾을 수 없습니다: {document_path}")
            # docs 폴더 내 파일 목록 제안
            docs_folder = self.settings["paths"]["docs_folder"]
            if os.path.exists(docs_folder):
                files = [f for f in os.listdir(docs_folder) if os.path.isfile(os.path.join(docs_folder, f))]
                if files:
                    print(f"💡 '{docs_folder}' 폴더의 사용 가능한 파일:")
                    for file in files[:10]:  # 최대 10개만 표시
                        print(f"   - {file}")
                    if len(files) > 10:
                        print(f"   ... 외 {len(files) - 10}개 파일")
            return None
        
        # 컬렉션명이 제공되지 않으면 자동 생성
        if collection_name is None:
            collection_name = self.generate_collection_name(full_path)
            print(f"📝 자동 생성된 컬렉션명: {collection_name}")
        
        return self.create_collection(collection_name, full_path)
    
    def get_collection_details(self, collection_name):
        """특정 컬렉션의 상세 정보 조회"""
        try:
            collection = self.client.get_collection(name=collection_name)
            count = collection.count()
            
            print(f"\n=== '{collection_name}' 컬렉션 상세 정보 ===")
            print(f"총 청크 수: {count}")
            
            if count > 0:
                # 샘플 데이터 가져오기
                sample = collection.peek(limit=5)
                
                if sample['metadatas'] and len(sample['metadatas']) > 0:
                    metadata = sample['metadatas'][0]
                    print(f"문서명: {metadata.get('document_name', 'Unknown')}")
                    print(f"파일경로: {metadata.get('document_path', 'Unknown')}")
                    print(f"인코딩: {metadata.get('encoding', 'Unknown')}")
                
                print(f"\n첫 {min(5, count)}개 청크 미리보기:")
                for i, (doc, meta) in enumerate(zip(sample['documents'], sample['metadatas'])):
                    if doc and meta:
                        preview = doc[:100] + "..." if len(doc) > 100 else doc
                        print(f"{i+1}. [{meta.get('chunk_id', i)}] {preview}")
            
            return collection
            
        except Exception as e:
            print(f"❌ 컬렉션 정보를 가져올 수 없습니다: {e}")
            return None
    
    def delete_collection(self, collection_name):
        """컬렉션 삭제 - ChromaDB 내부 파일도 함께 정리"""
        try:
            # 컬렉션이 존재하는지 확인
            existing_collections = [col.name for col in self.client.list_collections()]
            if collection_name not in existing_collections:
                print(f"❌ '{collection_name}' 컬렉션을 찾을 수 없습니다.")
                return False
            
            # 컬렉션 정보 먼저 가져오기
            collection = self.client.get_collection(name=collection_name)
            count = collection.count()
            
            # 삭제 확인
            print(f"⚠️  '{collection_name}' 컬렉션을 삭제하시겠습니까?")
            print(f"   - 총 청크 수: {count}개")
            response = input("이 작업은 되돌릴 수 없습니다. (y/N): ").strip().lower()
            if response not in ['y', 'yes']:
                print("삭제를 취소했습니다.")
                return False
            
            # 컬렉션 삭제
            print("🗑️ 컬렉션 및 관련 파일 삭제 중...")
            self.client.delete_collection(name=collection_name)
            
            # ChromaDB 디렉토리에서 관련 파일 정리 (선택적)
            import shutil
            import glob
            chroma_path = os.path.abspath(self.db_path)
            
            # 컬렉션 UUID 기반 디렉토리들 확인
            collection_dirs = glob.glob(os.path.join(chroma_path, "*-*-*-*-*"))
            
            # 빈 디렉토리나 임시 파일 정리
            cleaned_count = 0
            for dir_path in collection_dirs:
                try:
                    # 디렉토리가 비어있거나 불완전한 경우 정리
                    if os.path.isdir(dir_path):
                        files = os.listdir(dir_path)
                        # 필수 파일들이 없거나 비어있는 경우
                        required_files = ['data_level0.bin', 'header.bin', 'length.bin', 'link_lists.bin']
                        if not any(f in files for f in required_files) or len(files) == 0:
                            shutil.rmtree(dir_path)
                            cleaned_count += 1
                except Exception as e:
                    # 정리 중 오류가 발생해도 계속 진행
                    pass
            
            if cleaned_count > 0:
                print(f"🧹 정리된 빈 디렉토리: {cleaned_count}개")
            
            print(f"✅ '{collection_name}' 컬렉션이 완전히 삭제되었습니다.")
            return True
            
        except Exception as e:
            print(f"❌ 컬렉션 삭제 중 오류가 발생했습니다: {e}")
            return False
    
    def cleanup_database(self):
        """ChromaDB 데이터베이스 전체 정리"""
        try:
            print("🧹 ChromaDB 데이터베이스 정리 중...")
            
            # 현재 컬렉션 목록 확인
            collections = self.client.list_collections()
            active_collections = [col.name for col in collections]
            
            print(f"활성 컬렉션: {len(active_collections)}개")
            
            # ChromaDB 디렉토리 정리
            import shutil
            import glob
            chroma_path = os.path.abspath(self.db_path)
            
            if not os.path.exists(chroma_path):
                print("ChromaDB 디렉토리가 존재하지 않습니다.")
                return True
            
            # 컬렉션 UUID 기반 디렉토리들 확인
            collection_dirs = glob.glob(os.path.join(chroma_path, "*-*-*-*-*"))
            cleaned_count = 0
            
            for dir_path in collection_dirs:
                try:
                    if os.path.isdir(dir_path):
                        files = os.listdir(dir_path)
                        # 필수 파일들이 없거나 비어있는 디렉토리 정리
                        required_files = ['data_level0.bin', 'header.bin', 'length.bin', 'link_lists.bin']
                        if not any(f in files for f in required_files) or len(files) == 0:
                            shutil.rmtree(dir_path)
                            cleaned_count += 1
                            print(f"   🗑️ 정리됨: {os.path.basename(dir_path)}")
                except Exception as e:
                    print(f"   ⚠️ 정리 중 오류: {os.path.basename(dir_path)} - {e}")
                    continue
            
            print(f"✅ 데이터베이스 정리 완료 - 정리된 디렉토리: {cleaned_count}개")
            return True
            
        except Exception as e:
            print(f"❌ 데이터베이스 정리 중 오류: {e}")
            return False
    
    def show_help(self):
        """도움말 표시"""
        print("\n" + "="*60)
        print("📚 TinyRAG - 가벼운 오프라인 문서 검색 시스템")
        supported_extensions = self.get_supported_extensions()
        print(f"지원 형식: {', '.join(supported_extensions)}")
        print("\n명령어:")
        print("  add <파일경로> [컬렉션명]  : 새 문서 추가 (컬렉션명 생략시 자동생성)")
        print("  add <파일명>             : docs 폴더의 파일 추가 (파일명만으로도 가능)")
        print("  add \"파일 이름.docx\"      : 공백이 있는 파일명은 따옴표 사용")
        print("  list                     : 컬렉션 목록 보기")
        print("  detail <컬렉션명>        : 컬렉션 상세 정보 보기")
        print("  delete <컬렉션명>        : 컬렉션 삭제 (관련 파일도 정리)")
        print("  cleanup                  : 데이터베이스 전체 정리")
        print("  extensions               : 지원되는 파일 형식 상세 보기")
        print("  help                     : 이 도움말 표시")
        print("  quit                     : 종료")
        print("="*60)

def main():
    # 문서 처리기 초기화
    processor = DocumentProcessor()
    
    # 지원되는 파일 형식 표시
    supported_extensions = processor.get_supported_extensions()
    
    # 대화형 모드 (기본 문서 처리 제거)
    processor.show_help()
    
    while True:
        try:
            command = input("\n💾 명령: ").strip()
            
            if command.lower() in ['quit', 'exit', '종료', 'q']:
                print("👋 시스템을 종료합니다.")
                break
            
            elif command.lower() == 'list':
                processor.list_collections()
            
            elif command.lower() in ['help', 'h', '도움말']:
                processor.show_help()
            
            elif command.startswith('delete '):
                collection_name = command[7:].strip()
                if collection_name:
                    processor.delete_collection(collection_name)
                else:
                    print("사용법: delete <컬렉션명>")
            
            elif command.lower() == 'cleanup':
                print("⚠️ 데이터베이스 전체를 정리하시겠습니까?")
                response = input("사용하지 않는 파일들이 삭제됩니다. (y/N): ").strip().lower()
                if response in ['y', 'yes']:
                    processor.cleanup_database()
                else:
                    print("정리를 취소했습니다.")
            
            elif command.lower() == 'extensions':
                print("\n=== 지원되는 파일 형식 ===")
                extensions = processor.get_supported_extensions()
                
                print("📄 텍스트 파일:")
                print("  - .txt, .md (항상 지원)")
                
                if MARKITDOWN_AVAILABLE:
                    print("🚀 MarkItDown (Microsoft) - 모든 Office 문서 고품질 변환:")
                    print("  - .pdf, .docx, .doc, .xlsx, .xls, .pptx, .ppt")
                    print("  - Markdown 형식으로 최적화된 변환 제공")
                else:
                    print("🚀 MarkItDown: ❌ (pip install markitdown)")
                
                if DOCX_AVAILABLE:
                    print("📝 Microsoft Word (기본 지원):")
                    print("  - .docx, .doc (Markdown 형식으로 변환)")
                else:
                    print("📝 Microsoft Word: ❌ (python-docx 라이브러리 필요)")
                
                if PDF_AVAILABLE or PYMUPDF_AVAILABLE:
                    print("📑 PDF 파일 (기본 지원):")
                    if PYMUPDF_AVAILABLE:
                        print("  - .pdf (PyMuPDF 사용 - 고품질)")
                    else:
                        print("  - .pdf (PyPDF2 사용)")
                else:
                    print("📑 PDF 파일: ❌ (PyPDF2 또는 PyMuPDF 라이브러리 필요)")
                
                if EXCEL_AVAILABLE:
                    print("📊 Microsoft Excel (기본 지원):")
                    print("  - .xlsx, .xls (Markdown 표 형식으로 변환)")
                else:
                    print("📊 Microsoft Excel: ❌ (openpyxl 라이브러리 필요)")
                
                if PPTX_AVAILABLE:
                    print("📺 Microsoft PowerPoint (기본 지원):")
                    print("  - .pptx, .ppt (Markdown 형식으로 변환)")
                else:
                    print("📺 Microsoft PowerPoint: ❌ (python-pptx 라이브러리 필요)")
                
                print(f"\n총 지원 확장자: {', '.join(extensions)}")
                print("\n⭐ 권장 설치 (최고 품질):")
                if not MARKITDOWN_AVAILABLE:
                    print("  pip install markitdown")
                print("\n선택적 설치 (기본 지원용):")
                missing = []
                if not DOCX_AVAILABLE: missing.append("python-docx")
                if not (PDF_AVAILABLE or PYMUPDF_AVAILABLE): missing.append("PyPDF2 또는 PyMuPDF")
                if not EXCEL_AVAILABLE: missing.append("openpyxl")
                if not PPTX_AVAILABLE: missing.append("python-pptx")
                
                if missing:
                    print(f"  pip install {' '.join(missing)}")
                else:
                    print("  모든 기본 라이브러리가 설치되어 있습니다! 🎉")
            
            elif command.startswith('detail '):
                collection_name = command[7:].strip()
                if collection_name:
                    processor.get_collection_details(collection_name)
                else:
                    print("사용법: detail <컬렉션명>")
            
            elif command.startswith('add '):
                # 파일명에 공백이 있을 수 있으므로 더 정확한 파싱
                command_part = command[4:].strip()  # 'add ' 제거
                
                # 큰따옴표로 둘러싸인 파일 경로 처리
                if command_part.startswith('"'):
                    # "파일경로" 컬렉션명 형태
                    end_quote = command_part.find('"', 1)
                    if end_quote != -1:
                        file_path = command_part[1:end_quote]
                        remaining = command_part[end_quote+1:].strip()
                        collection_name = remaining if remaining else None
                    else:
                        print("❌ 파일 경로의 따옴표가 제대로 닫히지 않았습니다.")
                        continue
                else:
                    # 공백을 기준으로 분할하되, 마지막 부분을 컬렉션명으로 처리
                    parts = command_part.split()
                    if len(parts) == 0:
                        print("사용법: add <파일경로> [컬렉션명]")
                        continue
                    elif len(parts) == 1:
                        file_path = parts[0]
                        collection_name = None
                    else:
                        # 마지막 부분이 파일이 아니면 컬렉션명으로 간주
                        potential_file = ' '.join(parts[:-1])
                        potential_collection = parts[-1]
                        
                        # docs 폴더에서 파일 찾기 시도
                        potential_file_full = get_docs_path(potential_file)
                        full_path_all = get_docs_path(' '.join(parts))
                        
                        if os.path.exists(potential_file_full):
                            file_path = potential_file
                            collection_name = potential_collection
                        elif os.path.exists(full_path_all):
                            file_path = ' '.join(parts)
                            collection_name = None
                        else:
                            file_path = potential_file
                            collection_name = potential_collection
                
                # 문서 추가 시도
                processor.add_document(file_path, collection_name)
            
            elif command == '':
                continue
            
            else:
                print("알 수 없는 명령입니다. 'help'를 입력하여 사용 가능한 명령을 확인하세요.")
                
        except KeyboardInterrupt:
            print("\n\n👋 시스템을 종료합니다.")
            break
        except Exception as e:
            print(f"오류가 발생했습니다: {e}")

if __name__ == "__main__":
    main()
